from docx import Document
from pptx import Presentation
from pptx.util import Pt

def convert_word_to_ppt(word_file, ppt_file):
    # Load the Word document
    doc = Document(word_file)
    
    # Create a PowerPoint presentation
    prs = Presentation()
    
    # Iterate through the Word document paragraphs
    current_slide = None
    for para in doc.paragraphs:
        if para.style.name == 'Heading 1':
            # Create a new slide with Title layout for the first slide
            if current_slide is None:
                slide_layout = prs.slide_layouts[0]  # Title slide layout
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
            
            # Add slide
            current_slide = prs.slides.add_slide(slide_layout)
            
            # Set the title of the slide
            title = current_slide.shapes.title
            title.text = para.text
        else:
            if current_slide and para.text.strip():
                # Add content to the current slide
                content = current_slide.shapes.placeholders[1]
                text_frame = content.text_frame
                p = text_frame.add_paragraph()
                
                # Clear the default placeholder text
                if not text_frame.text:  # if text_frame is empty, use it for the first time
                    text_frame.text = ""

                # Add runs to preserve formatting
                for run in para.runs:
                    r = p.add_run()
                    r.text = run.text
                    # Apply bold
                    if run.bold:
                        r.font.bold = True
                    # Apply italic
                    if run.italic:
                        r.font.italic = True
                    # Apply underline
                    if run.underline:
                        r.font.underline = True
                    # Apply font size
                    if run.font.size:
                        r.font.size = Pt(run.font.size.pt)

    # Save the PowerPoint presentation
    prs.save(ppt_file)

# Example usage
convert_word_to_ppt('input.docx', 'output.pptx')
